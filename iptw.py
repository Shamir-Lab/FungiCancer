import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from scipy.stats import ttest_ind, mannwhitneyu
from statsmodels.stats.multitest import fdrcorrection
from statsmodels.stats.multitest import fdrcorrection
import math
from sklearn.linear_model import LogisticRegression
from statsmodels.formula.api import ols
import statsmodels.api as sm
import itertools
import matplotlib.pyplot as plt
import seaborn as sns
import matplotlib.pyplot as plt
#from matplotlib_venn import venn2, venn2_circles, venn3, venn3_circles
from scipy.stats import pearsonr
import plotting
from load_data import load_batch_corrected_files, add_bmi_and_obesity_cols
import utils
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
import os

# Apply FDR correction within each group using transform
def fdr_correction(p_values):
    return fdrcorrection(p_values)[1]

def apply_fdr(results):
    """
    #fdr correction for p-values
    p_values = results["p_value"].values
    mask = ~np.isnan(p_values)
    corrected_p_values = np.empty_like(p_values)
    corrected_p_values.fill(np.nan)
    corrected_p_values[mask] = fdrcorrection(p_values[mask])[1]
    results["p_value_fdr"] = corrected_p_values
    """
    results_no_na = results.dropna(subset = ["p_value"]).copy()  # Dropping nans - dropping cancers s.t cancer + var that don't have enough data (one or both groups i.e female + COAD : n<20).
    results_no_na['p_value_fdr'] = results_no_na.groupby('cancer_investigation')['p_value'].transform(fdr_correction)

    #print(results[results['p_value_fdr'] <0.05])
    return results_no_na

def IPTW_weighting(data, confounder_list, variable, need_dummies,list_of_bacteria, op1 = 1, op2 =0,  weight = "IPTW"):
    """
    confounder list - the list of confounders and variables of interest, can include the variable of interest, it will be dropped :)
    op1 and op2 are the 2 options for values in the variable column. if they are not 1 and 0, this will make sure to replace it
    need dummies - all vars that need one hot encoding
    list_of_bacteria - col names that contain counts
    """
    data= data.copy().replace({'FEMALE': 1, 'MALE': 0}) #makes life easier
    data = data.dropna(subset=confounder_list, axis=0) #make sure no null values in variable of interest + confounders
    data= data.replace({op1: 1, op2: 0}) #for the variable of interest - let's say 2 race groups which are usually one-hot encoded :)
    data = data.fillna(0)
    # create dummies for specified cols those columns
    data_with_dummies = pd.concat([data[confounder_list + ["investigation"]], pd.get_dummies(data[need_dummies])], axis = 1) #encoding
    data_with_dummies = data_with_dummies.drop(columns=need_dummies) #getting rid of non encoded columns

    columns_to_test = list_of_bacteria #fungi types
    results = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject'])

    for col in columns_to_test:
        for group in data["investigation"].drop_duplicates().values:
            #make sure both groups are large enough before the whole analysis:
            if (len(data_with_dummies[(data_with_dummies["investigation"] == group) & (data_with_dummies[variable] == 1)])>=20) and (len(data_with_dummies[(data_with_dummies["investigation"] == group) & (data_with_dummies[variable] == 0)])>=20) :
                X = data_with_dummies[data_with_dummies["investigation"] == group].drop([variable, "investigation"], axis=1).values #we don't want the dependant var in X, or the cancer
                y = data_with_dummies[data_with_dummies["investigation"] == group][variable].values
                model = LogisticRegression(solver='lbfgs', n_jobs=-1, class_weight='balanced') #took all parameters from Dan's example
                model.fit(X, y)

                #p_treated = float(np.sum(y==1))/y.shape[0]
                propensity_scores = model.predict_proba(X)[:, 1]
                df = data[data["investigation"] == group].copy()
                df["weights_for_%s_analysis"%variable] = 0
                if weight == "IPTW":
                    IP_treated = 1 / propensity_scores
                    IP_untreated = (1 / (1 - propensity_scores))
                elif weight == "MW":
                    IP_treated = min(propensity_scores, 1-propensity_scores) / propensity_scores
                    IP_untreated = min(propensity_scores, 1-propensity_scores) / (1-propensity_scores)
                df.loc[df[variable]==1, "weights_for_%s_analysis"%variable] = IP_treated[df[variable]==1]
                df.loc[df[variable]==0, "weights_for_%s_analysis"%variable] = IP_untreated[df[variable]==0]

                formula = '{col} ~ {variable}'.format(col=col, variable = variable)
                #print(formula)
                model = sm.WLS.from_formula(formula, data=df, weights=df["weights_for_%s_analysis"%variable]).fit()

                # Get the p-value for the x coefficient
                p_value = model.pvalues['%s'%variable]
            else:
                p_value = np.nan



            results.loc[len(results)] = [col, group, p_value, np.nan, False]
    results_new = apply_fdr(results.copy())
    return results_new, results


def analyze_iptw_flow(df, confounders, need_dummies, variable,taxa_cols ,op1, op2):
    df_res = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject', 'm_transform', 'm_bc'])

    filtered_df = df[df[variable].isin([op1, op2])].copy()
    results_new, results = IPTW_weighting(filtered_df, confounders, variable ,need_dummies, taxa_cols, op1, op2, weight="IPTW")
    #plotting.single_barplot_for_PSM(res, color_mapping)

    return results_new,results

def all_combos_batch_analysis(variable,batch_corr_dataframes,dataframes,fungi_cols, confounders, need_dummies, op1, op2):

    df_res_before = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject', 'm_transform', 'm_bc'])
    df_res = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject', 'm_transform', 'm_bc'])

    for key, df in batch_corr_dataframes.items():
        #print(key)
        method_transform, method_bc = key.rsplit('_', 1)
        df = add_bmi_and_obesity_cols(df, dataframes) #for bmi analysis - obese vs not obese
        res,res_before = analyze_iptw_flow(df, confounders, need_dummies, variable , fungi_cols ,op1, op2) #delete color_mapping from here
        res['m_transform'] = method_transform
        res_before['m_transform']= method_transform
        res['m_bc'] =  method_bc
        res_before['m_bc'] = method_bc

        # plotting.single_barplot_for_PSM(res, color_mapping,  f'Poore 2022-ALL FUNGI-{method_transform}-{method_bc}-{op1} vs. {op2} ONLY PRIMARY TUMOR')

        df_res = pd.concat([df_res, res])
        df_res_before = pd.concat([df_res_before, res_before])
    return df_res,df_res_before

def generate_p_values_per_batch_corrected_normalize_pair(comparisons_dict,dataframes,root_path,res_path,fungi_cols_WGS,fungi_cols):
    for specimen_type in ['RNA-Seq']:
        for cohort_definition in ['iptw']:
            if (((specimen_type == 'RNA-Seq')  & (cohort_definition in ['no_harvard','iptw'])) or ((specimen_type == 'WGS')  & (cohort_definition in ['all','iptw']))):
                print(f'{specimen_type}_{cohort_definition} - generation of p-values')
                batch_corr_dataframes = load_batch_corrected_files(dataframes, root_path,specimen_type,cohort_definition)

                for ind in range(len(comparisons_dict.keys())):
                    variable = list(comparisons_dict.keys())[ind]
                    factor = comparisons_dict[variable][4]

                    op1 = comparisons_dict[variable][0]
                    op2 = comparisons_dict[variable][1]
                    confounders = comparisons_dict[variable][2]
                    need_dummies = comparisons_dict[variable][3]

                    print(f'calculate p-values for {factor}')
                    if (specimen_type == 'WGS'):
                        df_res,df_res_before = all_combos_batch_analysis(factor,batch_corr_dataframes,dataframes,fungi_cols_WGS, confounders, need_dummies, op1, op2)# ,mul_type)
                    if (specimen_type == 'RNA-Seq'):
                        df_res,df_res_before = all_combos_batch_analysis(factor,batch_corr_dataframes,dataframes,fungi_cols, confounders, need_dummies, op1, op2)# ,mul_type)

                    file_name = (f'{specimen_type}_{cohort_definition}_{factor}_{op1}_{op2}.csv')
                    df_res_before.to_csv(f'{res_path}/{file_name}')

    return('Batch-corrected and normalized csv files were generated')
def apply_mul(results, mul_type):

    if (mul_type == 'fdr'):
      results_no_na = results.dropna(subset = ["p_value"]).copy()  # Dropping nans - dropping cancers s.t cancer + var that don't have enough data (one or both groups i.e female + COAD : n<20).
      results_no_na['p_value_fdr'] = results_no_na.groupby('cancer_investigation')['p_value'].transform(fdr_correction)
    if(mul_type == 'bon'):
      results_no_na = results.dropna(subset=["p_value"]).copy()  # Dropping nans - dropping cancers s.t cancer + var that don't have enough data (one or both groups i.e female + COAD : n<20).
      num_tests = results_no_na["cancer_investigation"].nunique()  # Number of tests conducted
      results_no_na['p_value_fdr'] = results_no_na['p_value'] * num_tests

    return results_no_na


def ppt_with_pivot_tables(specimen_type,cohort_definition,dataframes,root_path,res_path,comparisons_dict):
  print(f'{specimen_type} - {cohort_definition}')
  # set ppt
  prs = Presentation()
  prs.slide_width = Inches(13.333)  # 16:9 aspect ratio width in inches
  prs.slide_height = Inches(7.5)  # 16:9 aspect ratio height in inches
  slide_first = 0

  batch_corr_dataframes = load_batch_corrected_files(dataframes, root_path,specimen_type,cohort_definition)

  for ind in range(len(comparisons_dict.keys())):
    variable = list(comparisons_dict.keys())[ind]
    factor = comparisons_dict[variable][4]

    op1 = comparisons_dict[variable][0]
    op2 = comparisons_dict[variable][1]
    confounders = comparisons_dict[variable][2]
    need_dummies = comparisons_dict[variable][3]

    file_name = (f'{specimen_type}_{cohort_definition}_{factor}_{op1}_{op2}.csv')
    df_res_before = pd.read_csv(f'{res_path}/{file_name}', index_col=[0])

    df_mul = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject', 'm_transform', 'm_bc','mul_type'])
    print(factor)

    for key, df in batch_corr_dataframes.items():
        #print(key)
        method_transform, method_bc = key.rsplit('_', 1)
        temp = df_res_before[(df_res_before.m_transform == method_transform) & (df_res_before.m_bc ==method_bc)]

        for mul_type in ['fdr','bon']:
          temp2 = apply_mul(temp,mul_type)
          temp2['mul_type'] = mul_type
          df_mul = pd.concat([df_mul, temp2])


    title_slide_layout = prs.slide_layouts[slide_first]
    slide = prs.slides.add_slide(title_slide_layout)

    for mul_type in ['fdr','bon']:
      for p_val_threshold in [0.01, 0.05]:

        # generate pivot_table per multiple correction & p_values thresholds
        xxx = df_mul[df_mul.mul_type == mul_type]
        pivot_table = utils.pivot_fungi_counts(xxx,p_val_threshold)

        if len(pivot_table.columns)>0:
          taxas_per_table = utils.extract_significant_taxa(xxx, dataframes, max(list(pivot_table.columns)),p_val_threshold)
        else:
          taxas_per_table = utils.extract_significant_taxa(xxx, dataframes, 0,p_val_threshold)

        file_name = (f'{specimen_type}_{cohort_definition}_{factor}_{op1}_{op2}_{mul_type}_{p_val_threshold}.csv')
        pivot_table.to_csv(f'{res_path}/pivot_table_{file_name}')
        taxas_per_table.to_csv(f'{res_path}/taxas_per_table_{file_name}')

        #set title locatopn
        title_shape = slide.shapes.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(2)
        title_shape.text = f'{specimen_type} \n {factor} ({op1} vs. {op2})'
        title_shape.left =  Inches(cm_to_inches(4))  # Adjust left position
        title_shape.top = Inches(cm_to_inches(0.6))
        title_shape.width = Inches(cm_to_inches(25))  # Set a desired width
        title_shape.height = Inches(cm_to_inches(1))  # Adjust if necessary


        #shapes = slide.shapes
        tmp_df = pivot_table.reset_index()
        tmp_df = tmp_df.rename(columns={"cancer_investigation": "Cancer"})
        tmp_df = pd.DataFrame(np.vstack([tmp_df.columns, tmp_df]))

        cols = tmp_df.shape[1]
        rows = tmp_df.shape[0]
        if (mul_type == 'fdr'):
          left = Inches(cm_to_inches(0.1))
        else:
          left = Inches(cm_to_inches(18))
        #left = Inches(cm_to_inches(1))
        if (p_val_threshold == 0.01):
          top = Inches(cm_to_inches(3.5))
        else:
          top = Inches(cm_to_inches(10.5))
        width = Inches(cm_to_inches(1))
        height = Inches(cm_to_inches(0.8))


        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        txBox  = slide.shapes.add_textbox(left, top-Inches(cm_to_inches(1)), width, height)
        tf = txBox.text_frame
        tf.text = f'Multiple correction = {mul_type}, P-Val={p_val_threshold} =, Cohort={cohort_definition}'

        table.columns[0].width = Inches(cm_to_inches(2.5))

        for i in range(1,len(table.columns)):
          table.columns[i].width = Inches(cm_to_inches(1))

        for row_idx in range(0, len(table.rows)):
          for col_idx in range(0, len(table.columns)):
              value = tmp_df.iloc[row_idx,col_idx]
              cell = table.cell(row_idx, col_idx)
              cell.text = str(value)
              if (row_idx !=0):
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0xFF, 0xFA, 0xCD)

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)

    slide_first +=1

  prs.save(f'{res_path}/{specimen_type}_{cohort_definition}.pptx')

def extract_sig_species(specimen_type,cohort_definition,dataframes,root_path,res_path,comparisons_dict):
    print(f'{specimen_type} - {cohort_definition} - count tables of p-values')
    batch_corr_dataframes = load_batch_corrected_files(dataframes, root_path,specimen_type,cohort_definition)

    for ind in range(len(comparisons_dict.keys())):
        variable = list(comparisons_dict.keys())[ind]
        factor = comparisons_dict[variable][4]

        op1 = comparisons_dict[variable][0]
        op2 = comparisons_dict[variable][1]
        confounders = comparisons_dict[variable][2]
        need_dummies = comparisons_dict[variable][3]

        file_name = (f'{specimen_type}_{cohort_definition}_{factor}_{op1}_{op2}.csv')
        df_res_before = pd.read_csv(f'{res_path}/{file_name}', index_col=[0])

        df_mul = pd.DataFrame(columns=['fungus_id', 'cancer_investigation', 'p_value', 'p_value_fdr', 'reject', 'm_transform', 'm_bc','mul_type'])
        print(f'generate count tables of p-values for {factor}')

        for key, df in batch_corr_dataframes.items():
            #print(key)
            method_transform, method_bc = key.rsplit('_', 1)
            temp = df_res_before[(df_res_before.m_transform == method_transform) & (df_res_before.m_bc ==method_bc)]

            mul_type = 'bon'
            temp2 = apply_mul(temp,mul_type)
            temp2['mul_type'] = mul_type
            df_mul = pd.concat([df_mul, temp2])

        p_val_threshold = 0.05

        # generate pivot_table per multiple correction & p_values thresholds
        xxx = df_mul[df_mul.mul_type == mul_type]
        pivot_table = utils.pivot_fungi_counts(xxx,p_val_threshold)

        if len(pivot_table.columns)>0:
          taxas_per_table = utils.extract_significant_taxa(xxx, dataframes, max(list(pivot_table.columns)),p_val_threshold)
        else:
          taxas_per_table = utils.extract_significant_taxa(xxx, dataframes, 0,p_val_threshold)

        file_name = (f'{specimen_type}_{cohort_definition}_{factor}_{op1}_{op2}_{mul_type}_{p_val_threshold}.csv')
        pivot_table.to_csv(f'{res_path}/pivot_table_{file_name}')
        taxas_per_table.to_csv(f'{res_path}/taxas_per_table_{file_name}')

    return('p-values counts were extracted')

def calculate_mann_whitney(df, significant_fungi_df, dataframes,species_dict, is_salz):
    species_dict_id_to_genome = dict(zip(dataframes["poore_2022_fungi_species"]['genomeID'], dataframes["poore_2022_fungi_species"]['Species']))
    species_dict_id_to_genome = {key: value.replace('s__', '') for key, value in species_dict_id_to_genome.items()}

    results = []
    # Loop through significant fungi and perform the analysis
    for i, row in significant_fungi_df.iterrows():
        fungus = row["fungi_name"]
        if is_salz:
            if fungus == "Aspergillus pseudonomius":
                fungus = "Aspergillus pseudonomiae"
        else:
            fungus = species_dict.get(fungus, fungus)
        cancer_type = row["cancer"]
        factor = row["factor"]
        g1 = row["group1"]
        g2 = row["group2"]
        group1_values = df[(df["investigation"] == cancer_type) & (df[factor] == g1)][fungus]
        group2_values = df[(df["investigation"] == cancer_type) & (df[factor] == g2)][fungus]
        if group1_values.empty:
            reason = f"Insufficient data for group {g1}"
            p_value = None
            u = None
        elif group2_values.empty:
            reason = f"Insufficient data for group {g2}"
            p_value = None
            u = None
        else:
            u, p_value = mannwhitneyu(group1_values, group2_values, alternative='two-sided')
            reason = ''
        if is_salz:
            results.append({
                'fungus': fungus,
                'cancer_type': cancer_type.split("-")[1],
                'factor': factor,
                'group1': g1,
                'group2': g2,
                'p_value': format(p_value, '.3e') if p_value is not None else None,
                'size_g1': len(group1_values),
                'size_g2': len(group2_values),
                'mean_g1': round(group1_values.mean(), 3) if not group1_values.empty else None,
                'std_g1': round(group1_values.std(), 3) if not group1_values.empty else None,
                'mean_g2': round(group2_values.mean(), 3) if not group2_values.empty else None,
                'std_g2': round(group2_values.std(), 3) if not group2_values.empty else None,
                'non_zero_g1': (group1_values != 0).sum(),
                'non_zero_g2': (group2_values != 0).sum(),
                'status': 'Success' if reason == '' else 'Failed',
                'reason': reason
            })
        else:
            results.append({
                'fungus': species_dict_id_to_genome[fungus],
                'cancer_type': cancer_type.split("-")[1],
                'factor': factor,
                'group1': g1,
                'group2': g2,
                'p_value': format(p_value, '.3e') if p_value is not None else None,
                'size_g1': len(group1_values),
                'size_g2': len(group2_values),
                'mean_g1': round(group1_values.mean(), 3) if not group1_values.empty else None,
                'std_g1': round(group1_values.std(), 3) if not group1_values.empty else None,
                'mean_g2': round(group2_values.mean(), 3) if not group2_values.empty else None,
                'std_g2': round(group2_values.std(), 3) if not group2_values.empty else None,
                'non_zero_g1': (group1_values != 0).sum(),
                'non_zero_g2': (group2_values != 0).sum(),
                'status': 'Success' if reason == '' else 'Failed',
                'reason': reason
            })


    results_df = pd.DataFrame(results)
    return results_df



def create_summary_table_with_merged_cells(mann_whitney_df):
    df_copy = mann_whitney_df.copy()

    df_copy.rename(columns={
        'fungus': 'Specie',
        'cancer_type': 'Cancer',
        'factor': 'Factor',
        'group1': 'Group 1 Variable',
        'group2': 'Group 2 Variable',
        'size_g1': 'Group 1 N',
        'mean_g1': 'Group 1 Mean',
        'std_g1': 'Group 1 STD',
        'size_g2': 'Group 2 N',
        'mean_g2': 'Group 2 Mean',
        'std_g2': 'Group 2 STD',
        'p_value': 'P-Value',
        'non_zero_g1': 'Group 1 Non-Zero Count',
        'non_zero_g2': 'Group 2 Non-Zero Count'
    }, inplace=True)

    # Add columns for Mean±STD
    df_copy['Group 1 Mean±STD'] = df_copy.apply(
        lambda row: f"{row['Group 1 Mean']}±{row['Group 1 STD']}", axis=1
    )
    df_copy['Group 2 Mean±STD'] = df_copy.apply(
        lambda row: f"{row['Group 2 Mean']}±{row['Group 2 STD']}", axis=1
    )

    df_copy['Factor'] = df_copy['Factor'].replace({'race': 'Race', 'over_65': 'Age', 'gender': 'Gender', "obese":"BMI"})
    df_copy['Group 1 Variable'] = df_copy['Group 1 Variable'].replace({
        'WHITE': 'European', 'ASIAN': 'Asian', 'BLACK OR AFRICAN AMERICAN': 'African'
    })
    df_copy['Group 2 Variable'] = df_copy['Group 2 Variable'].replace({
        'WHITE': 'European', 'ASIAN': 'Asian', 'BLACK OR AFRICAN AMERICAN': 'African'
    })

    def replace_age_bmi(row):
        if row['Factor'] == 'Age':
            row['Group 1 Variable'] = '>70' if row['Group 1 Variable'] == 1 else '<70'
            row['Group 2 Variable'] = '>70' if row['Group 2 Variable'] == 1 else '<70'
        elif row['Factor'] == 'BMI':
            row['Group 1 Variable'] = '>30' if row['Group 1 Variable'] == 1 else '<30'
            row['Group 2 Variable'] = '>30' if row['Group 2 Variable'] == 1 else '<30'
        return row

    df_copy = df_copy.apply(replace_age_bmi, axis=1)

    df_copy['Group 1 Variable'] = df_copy['Group 1 Variable'].replace({'MALE': 'Male', 'FEMALE': 'Female'})
    df_copy['Group 2 Variable'] = df_copy['Group 2 Variable'].replace({'MALE': 'Male', 'FEMALE': 'Female'})



    # Define the column structure including Non-Zero Counts
    columns = [
        ('', 'Specie'),
        ('', 'Cancer'),
        ('', 'Factor'),
        ('Group 1', 'Variable'),
        ('Group 1', 'N'),
        ('Group 1', 'Mean±STD'),
        ('Group 1', 'Non-Zero Count'),
        ('Group 2', 'Variable'),
        ('Group 2', 'N'),
        ('Group 2', 'Mean±STD'),
        ('Group 2', 'Non-Zero Count'),
        ('', 'P-Value')
    ]

    # Reassign the column names with the new structure
    df_copy = df_copy[[
        'Specie', 'Cancer', 'Factor',
        'Group 1 Variable', 'Group 1 N', 'Group 1 Mean±STD', 'Group 1 Non-Zero Count',
        'Group 2 Variable', 'Group 2 N', 'Group 2 Mean±STD', 'Group 2 Non-Zero Count',
        'P-Value'
    ]]
    df_copy.columns = pd.MultiIndex.from_tuples(columns)

    return df_copy


def save_to_excel_with_merged_cells(mann_whitney_df, filename):
    wb = Workbook()
    ws = wb.active

    headers = [
        ['', '', '', 'Group 1', 'Group 1', 'Group 1', 'Group 1', 'Group 2', 'Group 2', 'Group 2', 'Group 2', ''],
        ['Specie', 'Cancer', 'Factor', 'Variable', 'N', 'Mean±STD', 'Non-Zero Count', 'Variable', 'N', 'Mean±STD', 'Non-Zero Count', 'P-Value']
    ]

    # Append header rows to the worksheet
    for header_row in headers:
        ws.append(header_row)

    # Add the data rows
    for r_idx, row in enumerate(dataframe_to_rows(mann_whitney_df, index=False, header=False), 3):
        ws.append(row)
        for c_idx, cell_value in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx)
            cell.number_format = '0.000' if isinstance(cell_value, float) else 'General'

    # Set values and formatting for individual header cells
    for col_num, header in enumerate(headers[1], 1):
        cell = ws.cell(row=2, column=col_num)
        cell.value = header
        cell.font = Font(bold=True)

    # Merge cells for Specie, Cancer, Factor, and P-Value
    columns_to_merge = [1, 2, 3, 12]  # Corrected the final index for P-Value column
    for col_num in columns_to_merge:
        col_letter = get_column_letter(col_num)
        ws.merge_cells(f'{col_letter}1:{col_letter}2')
        cell = ws[f'{col_letter}1']
        cell.value = headers[1][col_num - 1]
        cell.font = Font(bold=True)

    # Merge cells for Group 1 and Group 2 headers, now covering four columns each due to the additional 'Non-Zero Count' column
    ws.merge_cells('D1:G1')  # Corrected merge for Group 1
    ws['D1'].value = 'Group 1'
    ws['D1'].font = Font(bold=True)

    ws.merge_cells('H1:K1')  # Corrected merge for Group 2
    ws['H1'].value = 'Group 2'
    ws['H1'].font = Font(bold=True)

    # Save the workbook
    wb.save(filename)

def extract_significant_fungi(res_path):
    significant_fungi_data = []
    for file_name in os.listdir(res_path):
        if file_name.startswith('taxas_per_table') and file_name.endswith('.csv'):
            parts = file_name.replace('.csv', '').split('_')
            specimen_type = parts[3]  # RNA-Seq
            cohort_definition = parts[4]  # iptw
            if parts[5] == 'over':
                factor = f"{parts[5]}_{parts[6]}"  # Combine 'over' and '65'
                op1 = parts[7]  # group1 (e.g., 1)
                op2 = parts[8]  # group2 (e.g., 0)
            else:
                factor = parts[5]  # race or gender
                op1 = parts[6]  # group1 (e.g., WHITE)
                op2 = parts[7]  # group2 (e.g., BLACK OR AFRICAN AMERICAN)
            
            file_path = os.path.join(res_path, file_name)
            df = pd.read_csv(file_path)

            try:
                df = pd.read_csv(file_path)

                if df.empty:
                    print(f"File {file_name} is empty. Skipping.")
                    continue

                # Extract data from the CSV file
                for _, row in df.iterrows():
                    cancer_type = row['cancer_investigation']
                    fungus = row['fungus_id']
                    
                    significant_fungi_data.append({
                        'fungi_name': fungus,
                        'cancer': cancer_type,
                        'factor': factor,
                        'group1': op1,
                        'group2': op2
                    })
            
            except pd.errors.EmptyDataError:
                print(f"File {file_name} contains no data. Skipping.")
    
    
    return pd.DataFrame(significant_fungi_data)
